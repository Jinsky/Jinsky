"""
Klinik Merpati - Sistem Pakar Diagnosa Penyakit Avian
PowerPoint Presentation Generator
Design: Midnight Coral (Medical/Clinical Elegant)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ============================================================
# PRESENTATION SETUP
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ============================================================
# COLOR PALETTE - Midnight Coral (Medical/Clinical)
# ============================================================
C_PRIMARY = RGBColor(0x0F, 0x4C, 0x5C)        # Deep teal (dominance)
C_PRIMARY_DARK = RGBColor(0x09, 0x35, 0x42)   # Darker teal
C_SECONDARY = RGBColor(0xE8, 0xD8, 0xC4)      # Warm cream
C_ACCENT = RGBColor(0xE7, 0x6F, 0x51)         # Coral
C_GOLD = RGBColor(0xE9, 0xC4, 0x6A)           # Gold
C_MINT = RGBColor(0x5E, 0xB8, 0x9A)           # Mint
C_LIGHT_BG = RGBColor(0xF7, 0xF3, 0xEC)       # Cream background
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x1B, 0x26, 0x3B)           # Dark navy text
C_MUTED = RGBColor(0x5C, 0x6B, 0x7A)          # Muted text
C_LINE = RGBColor(0xD6, 0xC9, 0xB3)           # Soft line
C_CARD_BG = RGBColor(0xFB, 0xF7, 0xF0)        # Card background
C_CARD_BG_ALT = RGBColor(0xFF, 0xFC, 0xF5)    # Alternate card
C_RED = RGBColor(0xC1, 0x4B, 0x3A)
C_NAVY = RGBColor(0x1B, 0x2A, 0x4E)

# ============================================================
# PATHS
# ============================================================
BASE_DIR = r"C:\xampp\htdocs\Jinsky"
DIAGRAM_DIR = os.path.join(BASE_DIR, "diagram")
OUTPUT = os.path.join(BASE_DIR, "Klinik_Merpati_Presentasi_v2.pptx")

# ============================================================
# UTILITIES
# ============================================================

def add_text(slide, left, top, width, height, text, font_size=14,
             font_name="Calibri", bold=False, italic=False,
             color=C_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, space_after=0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, left, top, width, height, runs_data,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
             space_after=0, default_size=14, default_color=C_DARK):
    """runs_data: list of paragraphs where each paragraph is list of (text, opts)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(runs_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for txt, opts in para_runs:
            r = p.add_run()
            r.text = txt
            r.font.name = opts.get("font", "Calibri")
            r.font.size = Pt(opts.get("size", default_size))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", default_color)
    return tb


def add_rect(slide, left, top, width, height, fill=None, line=None,
             line_width=0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.shadow.inherit = False
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is not None:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    s.text_frame.text = ""
    return s


def add_oval(slide, left, top, width, height, fill=None, line=None,
             line_width=0, shadow=False):
    return add_rect(slide, left, top, width, height, fill, line, line_width,
                    MSO_SHAPE.OVAL, shadow)


def add_line(slide, x1, y1, x2, y2, color=C_PRIMARY, width=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def set_solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def new_slide(bg=C_WHITE):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, H, fill=bg)
    return slide


def add_slide_number(slide, num, total=15):
    add_text(slide, W - Inches(1.4), H - Inches(0.5), Inches(1.0), Inches(0.3),
             f"{num:02d} / {total:02d}",
             font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT,
             font_name="Calibri", bold=True)
    add_text(slide, Inches(0.5), H - Inches(0.5), Inches(8), Inches(0.3),
             "Klinik Merpati  •  Sistem Pakar Diagnosa Avian",
             font_size=10, color=C_MUTED, align=PP_ALIGN.LEFT,
             font_name="Calibri", bold=True)


def add_header_band(slide, kicker, title, kicker_color=C_ACCENT):
    """Standard content slide header with kicker + title"""
    # Small accent square
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.18),
             fill=kicker_color, shape=MSO_SHAPE.RECTANGLE)
    add_text(slide, Inches(0.85), Inches(0.45), Inches(8), Inches(0.35),
             kicker.upper(), font_size=11, color=kicker_color,
             font_name="Calibri", bold=True)
    add_text(slide, Inches(0.6), Inches(0.85), Inches(11.5), Inches(0.7),
             title, font_size=32, color=C_PRIMARY,
             font_name="Georgia", bold=True)


def add_footer_rule(slide, color=C_PRIMARY):
    add_rect(slide, 0, H - Inches(0.05), W, Inches(0.05), fill=color)


# ============================================================
# SLIDE 1: COVER
# ============================================================
def slide_01_cover():
    s = new_slide(bg=C_PRIMARY_DARK)

    # Decorative left side panel
    add_rect(s, 0, 0, Inches(5.2), H, fill=C_PRIMARY)

    # Decorative circles
    add_oval(s, Inches(-2.5), Inches(-2.5), Inches(5), Inches(5),
             fill=C_PRIMARY, line=None)
    add_oval(s, Inches(-1.5), Inches(5.5), Inches(3.5), Inches(3.5),
             fill=C_PRIMARY, line=None)
    add_oval(s, Inches(3.5), Inches(6), Inches(2.5), Inches(2.5),
             fill=C_PRIMARY, line=None)
    add_oval(s, Inches(-1.5), Inches(1.8), Inches(1.0), Inches(1.0),
             fill=C_GOLD, line=None)
    add_oval(s, Inches(4.0), Inches(2.5), Inches(0.6), Inches(0.6),
             fill=C_ACCENT, line=None)

    # Logo mark - bird silhouette in circle
    add_oval(s, Inches(0.6), Inches(0.6), Inches(0.7), Inches(0.7),
             fill=C_GOLD, line=None)
    add_text(s, Inches(0.6), Inches(0.6), Inches(0.7), Inches(0.7),
             "M", font_size=28, color=C_PRIMARY_DARK, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font_name="Georgia", bold=True)

    # Top tag
    add_text(s, Inches(1.4), Inches(0.7), Inches(3.5), Inches(0.5),
             "KLINIK MERPATI", font_size=12, color=C_LIGHT_BG,
             font_name="Calibri", bold=True)
    add_text(s, Inches(1.4), Inches(0.95), Inches(3.5), Inches(0.3),
             "Avian Diagnostic Expert System", font_size=9,
             color=C_SECONDARY, font_name="Calibri", italic=True)

    # Center content - bird icon symbol
    add_oval(s, Inches(1.6), Inches(2.6), Inches(2.0), Inches(2.0),
             fill=C_LIGHT_BG, line=None)
    add_oval(s, Inches(1.85), Inches(2.85), Inches(1.5), Inches(1.5),
             fill=C_PRIMARY, line=None)
    add_text(s, Inches(1.85), Inches(2.85), Inches(1.5), Inches(1.5),
             "+", font_size=70, color=C_GOLD, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font_name="Georgia", bold=True)

    # Bottom title block
    add_text(s, Inches(0.6), Inches(5.0), Inches(4.4), Inches(0.4),
             "SISTEM PAKAR", font_size=11, color=C_GOLD,
             font_name="Calibri", bold=True)
    add_text(s, Inches(0.6), Inches(5.35), Inches(4.4), Inches(1.7),
             "Diagnosa\nPenyakit\nAvian", font_size=28, color=C_LIGHT_BG,
             font_name="Georgia", bold=True, line_spacing=1.05)
    add_rect(s, Inches(0.6), Inches(6.85), Inches(0.6), Inches(0.05),
             fill=C_ACCENT, shape=MSO_SHAPE.RECTANGLE)

    # Right side content
    add_text(s, Inches(5.8), Inches(0.8), Inches(7), Inches(0.4),
             "SKRIPSI / TUGAS AKHIR", font_size=11, color=C_ACCENT,
             font_name="Calibri", bold=True)
    add_rect(s, Inches(5.8), Inches(1.15), Inches(0.4), Inches(0.04),
             fill=C_ACCENT, shape=MSO_SHAPE.RECTANGLE)

    add_text(s, Inches(5.8), Inches(1.4), Inches(7.2), Inches(2.6),
             "Klinik Merpati", font_size=64, color=C_LIGHT_BG,
             font_name="Georgia", bold=True, line_spacing=1.0)
    add_text(s, Inches(5.8), Inches(3.1), Inches(7.2), Inches(2.0),
             "Sistem Pakar Diagnosa\nPenyakit Avian dengan\nMetode Forward Chaining",
             font_size=26, color=C_SECONDARY, font_name="Georgia",
             italic=True, line_spacing=1.15)

    # Decorative right side
    add_rect(s, Inches(11.0), Inches(1.4), Inches(1.7), Inches(0.04),
             fill=C_GOLD, shape=MSO_SHAPE.RECTANGLE)

    # Bottom info block
    add_rect(s, Inches(5.8), Inches(5.4), Inches(6.9), Inches(0.04),
             fill=C_MUTED, shape=MSO_SHAPE.RECTANGLE)

    add_text(s, Inches(5.8), Inches(5.55), Inches(3.5), Inches(0.3),
             "DISUSUN OLEH", font_size=10, color=C_MUTED,
             font_name="Calibri", bold=True)
    add_text(s, Inches(5.8), Inches(5.85), Inches(3.5), Inches(0.4),
             "Nama Mahasiswa", font_size=15, color=C_LIGHT_BG,
             font_name="Georgia", bold=True)
    add_text(s, Inches(5.8), Inches(6.15), Inches(3.5), Inches(0.3),
             "NIM. 123456789", font_size=11, color=C_SECONDARY,
             font_name="Calibri")

    add_text(s, Inches(9.5), Inches(5.55), Inches(3.2), Inches(0.3),
             "PEMBIMBING", font_size=10, color=C_MUTED,
             font_name="Calibri", bold=True)
    add_text(s, Inches(9.5), Inches(5.85), Inches(3.2), Inches(0.4),
             "Nama Dosen, M.Kom", font_size=15, color=C_LIGHT_BG,
             font_name="Georgia", bold=True)
    add_text(s, Inches(9.5), Inches(6.15), Inches(3.2), Inches(0.3),
             "Program Studi Sistem Informasi", font_size=11, color=C_SECONDARY,
             font_name="Calibri")

    # Bottom year
    add_text(s, Inches(5.8), Inches(6.85), Inches(6.9), Inches(0.4),
             "Tahun 2025  •  Universitas Anda", font_size=11,
             color=C_MUTED, font_name="Calibri", bold=True)


# ============================================================
# SLIDE 2: LATAR BELAKANG
# ============================================================
def slide_02_latar_belakang():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB I  •  Pendahuluan", "Latar Belakang Penelitian")

    # Intro text
    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.8),
             "Merpati merupakan unggas yang rentan terhadap berbagai penyakit. "
             "Keterlambatan penanganan sering berujung pada kematian karena "
             "gejala awal sulit dikenali oleh peternak awam.",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    # Three big stat cards
    cards = [
        ("30", "Gejala", "tercatat dalam basis pengetahuan sistem", C_ACCENT),
        ("10", "Penyakit", "Avian yang dapat didiagnosa secara akurat", C_GOLD),
        ("92%", "Tingkat Akurasi", "berdasarkan metode Forward Chaining", C_MINT),
    ]
    card_w = Inches(3.85)
    card_h = Inches(2.0)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    y = Inches(2.8)
    for i, (num, label, desc, color) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y, card_w, card_h, fill=C_CARD_BG)
        add_rect(s, x, y, Inches(0.08), card_h, fill=color)
        add_text(s, x + Inches(0.35), y + Inches(0.25), card_w - Inches(0.5),
                 Inches(0.9), num, font_size=58, color=color,
                 font_name="Georgia", bold=True, line_spacing=1.0)
        add_text(s, x + Inches(0.35), y + Inches(1.15), card_w - Inches(0.5),
                 Inches(0.4), label, font_size=15, color=C_PRIMARY,
                 font_name="Calibri", bold=True)
        add_text(s, x + Inches(0.35), y + Inches(1.5), card_w - Inches(0.5),
                 Inches(0.5), desc, font_size=11, color=C_MUTED,
                 font_name="Calibri", line_spacing=1.25)

    # Problem context card
    add_rect(s, Inches(0.6), Inches(5.05), Inches(12.1), Inches(1.85),
             fill=C_PRIMARY)
    add_rect(s, Inches(0.6), Inches(5.05), Inches(0.08), Inches(1.85),
             fill=C_GOLD)

    add_text(s, Inches(0.95), Inches(5.25), Inches(11.5), Inches(0.4),
             "PERMASALAHAN UTAMA", font_size=11, color=C_GOLD,
             font_name="Calibri", bold=True)
    add_text(s, Inches(0.95), Inches(5.6), Inches(11.5), Inches(1.3),
             "Banyak peternak merpati belum memiliki akses cepat terhadap "
             "dokter hewan spesialis unggas, sehingga gejala awal sering "
             "terabaikan. Diperlukan sebuah sistem pakar yang dapat "
             "meniru penalaran ahli untuk memberikan diagnosa awal "
             "berbasis gejala secara mandiri, cepat, dan konsisten.",
             font_size=13, color=C_LIGHT_BG, font_name="Calibri",
             line_spacing=1.4)

    add_slide_number(s, 2)


# ============================================================
# SLIDE 3: RUMUSAN MASALAH
# ============================================================
def slide_03_rumusan_masalah():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB I  •  Pendahuluan", "Rumusan Masalah")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
             "Berdasarkan latar belakang di atas, permasalahan yang diangkat "
             "dalam penelitian ini adalah:",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    problems = [
        ("01", "Bagaimana merancang sistem pakar yang dapat mendiagnosa "
               "penyakit avian pada merpati berdasarkan gejala yang diamati?"),
        ("02", "Bagaimana menerapkan metode Forward Chaining sebagai "
               "mesin inferensi dalam proses diagnosa penyakit avian?"),
        ("03", "Bagaimana tingkat akurasi sistem pakar dalam memberikan "
               "hasil diagnosa penyakit avian pada merpati?"),
    ]

    y = Inches(2.5)
    for i, (num, text) in enumerate(problems):
        y_pos = y + Inches(1.05) * i
        # Number badge
        add_oval(s, Inches(0.7), y_pos + Inches(0.15), Inches(0.75),
                 Inches(0.75), fill=C_PRIMARY)
        add_text(s, Inches(0.7), y_pos + Inches(0.15), Inches(0.75),
                 Inches(0.75), num, font_size=18, color=C_GOLD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font_name="Georgia", bold=True)
        # Card
        add_rect(s, Inches(1.65), y_pos, Inches(11.05), Inches(0.95),
                 fill=C_CARD_BG)
        add_rect(s, Inches(1.65), y_pos, Inches(0.06), Inches(0.95),
                 fill=C_ACCENT)
        add_text(s, Inches(1.95), y_pos + Inches(0.18), Inches(10.5),
                 Inches(0.7), text, font_size=14, color=C_DARK,
                 font_name="Calibri", line_spacing=1.35)

    add_slide_number(s, 3)


# ============================================================
# SLIDE 4: BATASAN MASALAH
# ============================================================
def slide_04_batasan_masalah():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB I  •  Pendahuluan", "Batasan Masalah")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
             "Agar penelitian ini terfokus, maka batasan masalah "
             "yang ditetapkan adalah:",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    # 2x3 grid of constraint cards
    items = [
        ("Subjek Penelitian", "Merpati (Columba livia) sebagai objek studi utama.", C_ACCENT),
        ("Cakupan Penyakit", "10 jenis penyakit avian yang terdata pada basis pengetahuan.", C_GOLD),
        ("Metode Inferensi", "Forward Chaining berbasis aturan (rule-based).", C_MINT),
        ("Gejala", "30 gejala klinis yang menjadi dasar diagnosa.", C_PRIMARY),
        ("Platform", "Aplikasi berbasis web menggunakan PHP & MySQL.", C_NAVY),
        ("Pengguna", "Peternak merpati, peneliti, dan dokter hewan.", C_RED),
    ]
    card_w = Inches(3.95)
    card_h = Inches(1.7)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    start_x = Inches(0.6)
    start_y = Inches(2.5)
    for i, (title, desc, color) in enumerate(items):
        col = i % 3
        row = i // 3
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, fill=C_CARD_BG)
        # Icon circle
        add_oval(s, x + Inches(0.3), y + Inches(0.3), Inches(0.55),
                 Inches(0.55), fill=color)
        # Icon symbol
        symbols = ["✓", "⊕", "→", "◇", "</>", "♥"]
        add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(0.55),
                 Inches(0.55), symbols[i], font_size=20, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font_name="Georgia", bold=True)
        add_text(s, x + Inches(1.0), y + Inches(0.3), card_w - Inches(1.1),
                 Inches(0.4), title, font_size=14, color=C_PRIMARY,
                 font_name="Calibri", bold=True)
        add_text(s, x + Inches(0.35), y + Inches(1.0), card_w - Inches(0.5),
                 Inches(0.7), desc, font_size=11, color=C_MUTED,
                 font_name="Calibri", line_spacing=1.3)

    add_slide_number(s, 4)


# ============================================================
# SLIDE 5: TUJUAN PENELITIAN
# ============================================================
def slide_05_tujuan():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB I  •  Pendahuluan", "Tujuan Penelitian")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
             "Tujuan yang ingin dicapai dalam penelitian ini adalah:",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    # 2x2 large goal cards with numbers
    goals = [
        ("01", "Membangun Sistem", "Merancang dan membangun sistem pakar "
         "berbasis web untuk diagnosa penyakit avian pada merpati.",
         C_PRIMARY, C_GOLD),
        ("02", "Menerapkan Metode", "Mengimplementasikan metode Forward "
         "Chaining sebagai mesin inferensi utama sistem pakar.",
         C_ACCENT, C_LIGHT_BG),
        ("03", "Menguji Akurasi", "Menguji tingkat akurasi dan validitas "
         "sistem menggunakan metode Black-Box Testing.",
         C_GOLD, C_PRIMARY_DARK),
        ("04", "Memberikan Solusi", "Menyediakan alat bantu diagnosa awal "
         "yang mudah diakses oleh peternak merpati.",
         C_MINT, C_PRIMARY_DARK),
    ]
    card_w = Inches(5.95)
    card_h = Inches(1.95)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.6)
    start_y = Inches(2.5)
    for i, (num, title, desc, color, txt_color) in enumerate(goals):
        col = i % 2
        row = i // 2
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, fill=color)
        # Number in light text
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(2),
                 Inches(1.0), num, font_size=50, color=txt_color,
                 font_name="Georgia", bold=True, line_spacing=1.0)
        add_text(s, x + Inches(2.2), y + Inches(0.3), card_w - Inches(2.4),
                 Inches(0.5), title, font_size=17, color=txt_color,
                 font_name="Calibri", bold=True)
        add_text(s, x + Inches(2.2), y + Inches(0.78), card_w - Inches(2.4),
                 Inches(1.0), desc, font_size=11, color=txt_color,
                 font_name="Calibri", line_spacing=1.3)
        # Small accent line
        add_rect(s, x + Inches(0.3), y + card_h - Inches(0.1),
                 Inches(0.5), Inches(0.04), fill=txt_color)

    add_slide_number(s, 5)


# ============================================================
# SLIDE 6: MANFAAT PENELITIAN
# ============================================================
def slide_06_manfaat():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB I  •  Pendahuluan", "Manfaat Penelitian")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
             "Penelitian ini diharapkan memberikan manfaat bagi berbagai "
             "pihak, antara lain:",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    # 4-column benefits with icons
    benefits = [
        ("Peternak", "Mendapatkan alat bantu diagnosa awal yang praktis, "
         "cepat, dan mudah digunakan kapan saja.", C_ACCENT, "P"),
        ("Akademisi", "Menambah referensi dan kontribusi ilmiah di "
         "bidang sistem pakar dan informatika veteriner.", C_GOLD, "A"),
        ("Dokter Hewan", "Membantu proses skrining awal dan dokumentasi "
         "riwayat kesehatan merpati pasien.", C_MINT, "D"),
        ("Pengembang", "Menjadi dasar untuk pengembangan sistem serupa "
         "pada jenis unggas atau hewan lain.", C_PRIMARY, "I"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(3.3)
    gap = Inches(0.18)
    start_x = Inches(0.6)
    y = Inches(2.5)
    for i, (title, desc, color, letter) in enumerate(benefits):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y, card_w, card_h, fill=C_CARD_BG)
        # Top color block
        add_rect(s, x, y, card_w, Inches(1.2), fill=color)
        # Letter icon
        add_oval(s, x + card_w/2 - Inches(0.45), y + Inches(0.32),
                 Inches(0.9), Inches(0.9), fill=C_LIGHT_BG)
        add_text(s, x + card_w/2 - Inches(0.45), y + Inches(0.32),
                 Inches(0.9), Inches(0.9), letter, font_size=30,
                 color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font_name="Georgia", bold=True)
        # Title and description
        add_text(s, x + Inches(0.25), y + Inches(1.45), card_w - Inches(0.5),
                 Inches(0.5), "UNTUK", font_size=9, color=C_MUTED,
                 font_name="Calibri", bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.25), y + Inches(1.7), card_w - Inches(0.5),
                 Inches(0.5), title, font_size=18, color=C_PRIMARY,
                 font_name="Calibri", bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x + card_w/2 - Inches(0.2), y + Inches(2.2),
                 Inches(0.4), Inches(0.04), fill=color)
        add_text(s, x + Inches(0.25), y + Inches(2.35), card_w - Inches(0.5),
                 Inches(0.9), desc, font_size=11, color=C_MUTED,
                 font_name="Calibri", line_spacing=1.4, align=PP_ALIGN.CENTER)

    # Bottom callout
    add_rect(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.85),
             fill=C_PRIMARY)
    add_text(s, Inches(0.85), Inches(6.15), Inches(0.7), Inches(0.7),
             "✦", font_size=22, color=C_GOLD, font_name="Georgia",
             anchor=MSO_ANCHOR.MIDDLE, bold=True)
    add_text(s, Inches(1.4), Inches(6.2), Inches(11.0), Inches(0.55),
             "Kontribusi nyata dalam transformasi digital dunia peternakan "
             "dan kesejahteraan unggas di Indonesia.",
             font_size=13, color=C_LIGHT_BG, font_name="Calibri",
             italic=True, line_spacing=1.3)

    add_slide_number(s, 6)


# ============================================================
# SLIDE 7: METODE PENGUMPULAN DATA
# ============================================================
def slide_07_metode_data():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB III  •  Metodologi", "Metode Pengumpulan Data")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
             "Data penelitian dikumpulkan melalui tiga metode utama berikut:",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    methods = [
        ("Wawancara",
         "Melakukan sesi tanya jawab dengan dokter hewan dan peternak "
         "merpati berpengalaman untuk menggali pengetahuan pakar tentang "
         "gejala, diagnosis, serta penanganan penyakit avian.",
         "01", C_ACCENT),
        ("Observasi",
         "Mengamati secara langsung kondisi merpati yang sakit pada "
         "kandang peternakan untuk mencatat pola gejala yang muncul "
         "dan memverifikasi basis pengetahuan sistem.",
         "02", C_GOLD),
        ("Studi Pustaka",
         "Mengkaji literatur ilmiah, jurnal, buku veteriner, dan "
         "dokumen resmi terkait penyakit avian serta teori sistem "
         "pakar dengan metode Forward Chaining.",
         "03", C_MINT),
    ]
    card_h = Inches(1.25)
    gap = Inches(0.18)
    y = Inches(2.45)
    for i, (title, desc, num, color) in enumerate(methods):
        y_pos = y + (card_h + gap) * i
        # Background card
        add_rect(s, Inches(0.6), y_pos, Inches(12.1), card_h, fill=C_CARD_BG)
        # Number side
        add_rect(s, Inches(0.6), y_pos, Inches(2.0), card_h, fill=color)
        add_text(s, Inches(0.6), y_pos, Inches(2.0), card_h, num,
                 font_size=50, color=C_LIGHT_BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name="Georgia", bold=True)
        # Title and description
        add_text(s, Inches(2.95), y_pos + Inches(0.18), Inches(9.6),
                 Inches(0.4), title, font_size=18, color=C_PRIMARY,
                 font_name="Georgia", bold=True)
        add_rect(s, Inches(2.95), y_pos + Inches(0.62), Inches(0.4),
                 Inches(0.04), fill=color)
        add_text(s, Inches(2.95), y_pos + Inches(0.7), Inches(9.6),
                 Inches(0.55), desc, font_size=11, color=C_MUTED,
                 font_name="Calibri", line_spacing=1.3)

    add_slide_number(s, 7)


# ============================================================
# SLIDE 8: METODE PENGEMBANGAN SISTEM (Waterfall)
# ============================================================
def slide_08_metode_pengembangan():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB III  •  Metodologi", "Metode Pengembangan Sistem")

    # Intro
    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
             "Sistem dikembangkan dengan model Waterfall dengan lima "
             "tahapan sistematis:",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    # 5-stage horizontal process
    stages = [
        ("1", "Analisis", "Studi kebutuhan & identifikasi masalah", C_ACCENT),
        ("2", "Desain", "Perancangan sistem, basis data, dan antarmuka", C_GOLD),
        ("3", "Implementasi", "Pengkodean dengan PHP & MySQL", C_MINT),
        ("4", "Pengujian", "Black-Box Testing untuk validasi", C_PRIMARY),
        ("5", "Pemeliharaan", "Deployment & perawatan sistem", C_NAVY),
    ]
    n = len(stages)
    box_w = Inches(2.15)
    box_h = Inches(2.4)
    gap = Inches(0.25)
    total = box_w * n + gap * (n - 1)
    start_x = (W - total) / 2
    y = Inches(2.55)
    for i, (num, title, desc, color) in enumerate(stages):
        x = start_x + (box_w + gap) * i
        # Card
        add_rect(s, x, y, box_w, box_h, fill=C_CARD_BG)
        # Top color block
        add_rect(s, x, y, box_w, Inches(0.6), fill=color)
        add_text(s, x, y, box_w, Inches(0.6), f"TAHAP {num}",
                 font_size=11, color=C_LIGHT_BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name="Calibri", bold=True)
        # Big number
        add_oval(s, x + box_w/2 - Inches(0.4), y + Inches(0.8),
                 Inches(0.8), Inches(0.8), fill=color)
        add_text(s, x + box_w/2 - Inches(0.4), y + Inches(0.8),
                 Inches(0.8), Inches(0.8), num, font_size=28, color=C_LIGHT_BG,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font_name="Georgia", bold=True)
        # Title
        add_text(s, x + Inches(0.1), y + Inches(1.7), box_w - Inches(0.2),
                 Inches(0.4), title, font_size=15, color=C_PRIMARY,
                 font_name="Calibri", bold=True, align=PP_ALIGN.CENTER)
        # Description
        add_text(s, x + Inches(0.15), y + Inches(2.05), box_w - Inches(0.3),
                 Inches(0.4), desc, font_size=9, color=C_MUTED,
                 font_name="Calibri", line_spacing=1.25, align=PP_ALIGN.CENTER)
        # Arrow between
        if i < n - 1:
            arr_x = x + box_w + Inches(0.02)
            arr_y = y + box_h/2 - Inches(0.12)
            add_text(s, arr_x, arr_y, Inches(0.2), Inches(0.25),
                     ">", font_size=16, color=C_PRIMARY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font_name="Georgia", bold=True)

    # Bottom callout: Forward Chaining
    add_rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.55),
             fill=C_PRIMARY)
    add_rect(s, Inches(0.6), Inches(5.4), Inches(0.1), Inches(1.55),
             fill=C_GOLD)
    add_text(s, Inches(0.95), Inches(5.55), Inches(11.5), Inches(0.4),
             "MESIN INFERENSI", font_size=11, color=C_GOLD,
             font_name="Calibri", bold=True)
    add_text(s, Inches(0.95), Inches(5.85), Inches(11.5), Inches(0.5),
             "Forward Chaining", font_size=22, color=C_LIGHT_BG,
             font_name="Georgia", bold=True, italic=True)
    add_text(s, Inches(0.95), Inches(6.4), Inches(11.5), Inches(0.5),
             "Pencocokan gejala (fakta) dengan aturan IF–THEN secara "
             "berurutan hingga ditemukan penyakit yang paling sesuai.",
             font_size=12, color=C_SECONDARY, font_name="Calibri",
             line_spacing=1.3)

    add_slide_number(s, 8)


# ============================================================
# SLIDE 9: PERANCANGAN (Flowchart Lama & Baru)
# ============================================================
def slide_09_perancangan():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB IV  •  Perancangan", "Flowchart Sistem Lama & Baru")

    # Two columns
    col_w = Inches(6.0)
    col_h = Inches(4.7)
    y = Inches(1.85)
    gap = Inches(0.2)
    start_x = Inches(0.55)

    # === LEFT: Sistem Lama ===
    x = start_x
    add_rect(s, x, y, col_w, col_h, fill=C_CARD_BG)
    add_rect(s, x, y, col_w, Inches(0.55), fill=C_RED)
    add_text(s, x + Inches(0.3), y, Inches(2), Inches(0.55), "BEFORE",
             font_size=11, color=C_LIGHT_BG, anchor=MSO_ANCHOR.MIDDLE,
             font_name="Calibri", bold=True)
    add_text(s, x, y, col_w, Inches(0.55), "Sistem Lama",
             font_size=18, color=C_LIGHT_BG, anchor=MSO_ANCHOR.MIDDLE,
             font_name="Georgia", bold=True, align=PP_ALIGN.CENTER)
    # Image - centered in the box
    img_h_target = Inches(3.6)
    img_y = y + Inches(0.7)
    try:
        s.shapes.add_picture(os.path.join(DIAGRAM_DIR, "flowchart/sistem_lama_tight.png"),
                              x + Inches(0.3), img_y, height=img_h_target)
    except Exception as e:
        add_rect(s, x + Inches(0.3), img_y, col_w - Inches(0.6), img_h_target,
                 fill=C_CARD_BG_ALT)
        add_text(s, x + Inches(0.3), img_y, col_w - Inches(0.6), img_h_target,
                 "Flowchart Lama", font_size=14, color=C_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # === RIGHT: Sistem Baru ===
    x = start_x + col_w + gap
    add_rect(s, x, y, col_w, col_h, fill=C_CARD_BG)
    add_rect(s, x, y, col_w, Inches(0.55), fill=C_PRIMARY)
    add_text(s, x + Inches(0.3), y, Inches(2), Inches(0.55), "AFTER",
             font_size=11, color=C_GOLD, anchor=MSO_ANCHOR.MIDDLE,
             font_name="Calibri", bold=True)
    add_text(s, x, y, col_w, Inches(0.55), "Sistem Baru (Klinik Merpati)",
             font_size=18, color=C_LIGHT_BG, anchor=MSO_ANCHOR.MIDDLE,
             font_name="Georgia", bold=True, align=PP_ALIGN.CENTER)
    img_y = y + Inches(0.7)
    try:
        s.shapes.add_picture(os.path.join(DIAGRAM_DIR, "flowchart/sistem_baru_tight.png"),
                              x + Inches(0.3), img_y, height=img_h_target)
    except Exception as e:
        add_rect(s, x + Inches(0.3), img_y, col_w - Inches(0.6), img_h_target,
                 fill=C_CARD_BG_ALT)
        add_text(s, x + Inches(0.3), img_y, col_w - Inches(0.6), img_h_target,
                 "Flowchart Baru", font_size=14, color=C_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom comparison strip
    strip_y = Inches(6.5)
    strip_h = Inches(0.4)
    add_rect(s, Inches(0.55), strip_y, Inches(6.0), strip_h, fill=C_CARD_BG)
    add_text(s, Inches(0.55), strip_y, Inches(6.0), strip_h,
             "✕  Diagnosa manual oleh peternak",
             font_size=11, color=C_DARK, font_name="Calibri",
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, italic=True)
    add_rect(s, Inches(6.75), strip_y, Inches(6.0), strip_h, fill=C_PRIMARY)
    add_text(s, Inches(6.75), strip_y, Inches(6.0), strip_h,
             "✓  Diagnosa otomatis Forward Chaining",
             font_size=11, color=C_LIGHT_BG, font_name="Calibri",
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, bold=True)

    add_slide_number(s, 9)


# ============================================================
# SLIDE 10: USE CASE
# ============================================================
def slide_10_use_case():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB IV  •  Perancangan", "Diagram Use Case")

    # Left: image
    img_x = Inches(0.6)
    img_y = Inches(1.8)
    img_w = Inches(7.0)
    img_h = Inches(5.1)
    add_rect(s, img_x, img_y, img_w, img_h, fill=C_CARD_BG)
    try:
        s.shapes.add_picture(os.path.join(DIAGRAM_DIR, "activity/use_case.png"),
                              img_x + Inches(0.2), img_y + Inches(0.2),
                              width=img_w - Inches(0.4),
                              height=img_h - Inches(0.4))
    except Exception as e:
        add_text(s, img_x, img_y, img_w, img_h, "Use Case Diagram",
                 font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Right: actors & key features
    right_x = Inches(7.85)
    right_w = Inches(4.9)

    add_text(s, right_x, Inches(1.8), right_w, Inches(0.4),
             "AKTOR & FITUR UTAMA", font_size=11, color=C_ACCENT,
             font_name="Calibri", bold=True)
    add_rect(s, right_x, Inches(2.15), Inches(0.5), Inches(0.04),
             fill=C_ACCENT)

    # Aktor 1
    add_rect(s, right_x, Inches(2.4), right_w, Inches(1.55), fill=C_CARD_BG)
    add_oval(s, right_x + Inches(0.2), Inches(2.55), Inches(0.55),
             Inches(0.55), fill=C_ACCENT)
    add_text(s, right_x + Inches(0.2), Inches(2.55), Inches(0.55),
             Inches(0.55), "U", font_size=20, color=C_LIGHT_BG,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font_name="Georgia", bold=True)
    add_text(s, right_x + Inches(0.9), Inches(2.5), right_w - Inches(1),
             Inches(0.4), "User (Pengunjung)", font_size=15, color=C_PRIMARY,
             font_name="Calibri", bold=True)
    add_text(s, right_x + Inches(0.9), Inches(2.85), right_w - Inches(1),
             Inches(1.0),
             "• Melihat katalog penyakit\n• Melakukan konsultasi gejala\n"
             "• Melihat hasil diagnosa & riwayat",
             font_size=10.5, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4)

    # Aktor 2
    add_rect(s, right_x, Inches(4.1), right_w, Inches(1.55), fill=C_CARD_BG)
    add_oval(s, right_x + Inches(0.2), Inches(4.25), Inches(0.55),
             Inches(0.55), fill=C_GOLD)
    add_text(s, right_x + Inches(0.2), Inches(4.25), Inches(0.55),
             Inches(0.55), "A", font_size=20, color=C_PRIMARY_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font_name="Georgia", bold=True)
    add_text(s, right_x + Inches(0.9), Inches(4.2), right_w - Inches(1),
             Inches(0.4), "Admin", font_size=15, color=C_PRIMARY,
             font_name="Calibri", bold=True)
    add_text(s, right_x + Inches(0.9), Inches(4.55), right_w - Inches(1),
             Inches(1.0),
             "• Mengelola data gejala\n• Mengelola data penyakit & aturan\n"
             "• Melihat statistik pengunjung",
             font_size=10.5, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4)

    # Use cases total
    add_rect(s, right_x, Inches(5.8), right_w, Inches(1.1), fill=C_PRIMARY)
    add_text(s, right_x + Inches(0.3), Inches(5.9), Inches(2),
             Inches(0.3), "TOTAL USE CASE", font_size=10, color=C_GOLD,
             font_name="Calibri", bold=True)
    add_text(s, right_x + Inches(0.3), Inches(6.15), Inches(2),
             Inches(0.6), "14+", font_size=32, color=C_LIGHT_BG,
             font_name="Georgia", bold=True, line_spacing=1.0)
    add_text(s, right_x + Inches(2.0), Inches(6.0), right_w - Inches(2.2),
             Inches(0.85),
             "Interaksi sistem yang mencakup konsultasi, manajemen basis "
             "pengetahuan, dan pelacakan riwayat.",
             font_size=10, color=C_SECONDARY, font_name="Calibri",
             line_spacing=1.3, italic=True)

    add_slide_number(s, 10)


# ============================================================
# SLIDE 11: ACTIVITY DIAGRAM (intinya)
# ============================================================
def slide_11_activity():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB IV  •  Perancangan", "Activity Diagram — Alur Konsultasi")

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.5),
             "Alur aktivitas utama pada proses konsultasi dan diagnosa "
             "sistem pakar:",
             font_size=14, color=C_MUTED, font_name="Calibri",
             line_spacing=1.4, italic=True)

    # Horizontal flow with arrows - 6 steps
    steps = [
        ("Mulai", "User membuka halaman konsultasi", C_ACCENT),
        ("Input Nama", "Masukkan identitas pemilik merpati", C_GOLD),
        ("Pilih Gejala", "Tandai gejala yang teramati", C_MINT),
        ("Kirim Data", "Sistem menerima input gejala", C_PRIMARY),
        ("Proses", "Mesin inferensi mencocokkan aturan", C_NAVY),
        ("Hasil", "Tampil diagnosa & tingkat keyakinan", C_ACCENT),
    ]
    n = len(steps)
    box_w = Inches(1.85)
    box_h = Inches(2.1)
    arrow_w = Inches(0.18)
    total = box_w * n + arrow_w * (n - 1)
    start_x = (W - total) / 2
    y = Inches(2.55)
    for i, (title, desc, color) in enumerate(steps):
        x = start_x + (box_w + arrow_w) * i
        # Hexagon-like via rounded rect
        add_rect(s, x, y, box_w, box_h, fill=C_CARD_BG)
        add_rect(s, x, y, box_w, Inches(0.5), fill=color)
        add_text(s, x, y, box_w, Inches(0.5), f"LANGKAH {i+1}",
                 font_size=9, color=C_LIGHT_BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name="Calibri", bold=True)
        # Big number circle
        add_oval(s, x + box_w/2 - Inches(0.3), y + Inches(0.65),
                 Inches(0.6), Inches(0.6), fill=color)
        add_text(s, x + box_w/2 - Inches(0.3), y + Inches(0.65),
                 Inches(0.6), Inches(0.6), str(i+1), font_size=18,
                 color=C_LIGHT_BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name="Georgia", bold=True)
        # Title and desc
        add_text(s, x + Inches(0.1), y + Inches(1.32), box_w - Inches(0.2),
                 Inches(0.3), title, font_size=12, color=C_PRIMARY,
                 font_name="Calibri", bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(1.55), box_w - Inches(0.2),
                 Inches(0.55), desc, font_size=8, color=C_MUTED,
                 font_name="Calibri", line_spacing=1.25, align=PP_ALIGN.CENTER)
        # Arrow
        if i < n - 1:
            arr_x = x + box_w + Inches(0.0)
            arr_y = y + box_h/2 - Inches(0.15)
            add_text(s, arr_x, arr_y, arrow_w, Inches(0.3),
                     "▶", font_size=14, color=C_PRIMARY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font_name="Calibri", bold=True)

    # Decision diamond explanation
    add_rect(s, Inches(0.6), Inches(5.1), Inches(7.4), Inches(1.8),
             fill=C_PRIMARY)
    add_text(s, Inches(0.95), Inches(5.25), Inches(7), Inches(0.4),
             "KEPUTUSAN SISTEM", font_size=11, color=C_GOLD,
             font_name="Calibri", bold=True)
    add_text(s, Inches(0.95), Inches(5.6), Inches(7), Inches(0.5),
             "Apakah ada aturan yang cocok?", font_size=18, color=C_LIGHT_BG,
             font_name="Georgia", bold=True, italic=True)
    add_text(s, Inches(0.95), Inches(6.15), Inches(7), Inches(0.7),
             "Jika YA → tampilkan penyakit, persentase, dan solusi.\n"
             "Jika TIDAK → tampilkan pesan 'Diagnosis Tidak Ditemukan'.",
             font_size=11, color=C_SECONDARY, font_name="Calibri",
             line_spacing=1.4)

    # Side: output
    add_rect(s, Inches(8.25), Inches(5.1), Inches(4.55), Inches(1.8),
             fill=C_CARD_BG)
    add_text(s, Inches(8.5), Inches(5.25), Inches(4), Inches(0.4),
             "OUTPUT SISTEM", font_size=11, color=C_ACCENT,
             font_name="Calibri", bold=True)
    add_text(s, Inches(8.5), Inches(5.55), Inches(4), Inches(0.4),
             "Nama Penyakit", font_size=15, color=C_PRIMARY,
             font_name="Calibri", bold=True)
    add_text(s, Inches(8.5), Inches(5.9), Inches(4), Inches(0.4),
             "Tingkat Kecocokan (%)", font_size=11, color=C_DARK,
             font_name="Calibri")
    add_text(s, Inches(8.5), Inches(6.2), Inches(4), Inches(0.4),
             "Solusi & Pencegahan", font_size=11, color=C_DARK,
             font_name="Calibri")
    add_text(s, Inches(8.5), Inches(6.5), Inches(4), Inches(0.4),
             "Tersimpan ke Riwayat", font_size=11, color=C_DARK,
             font_name="Calibri")

    add_slide_number(s, 11)


# ============================================================
# SLIDE 12: ERD
# ============================================================
def slide_12_erd():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB IV  •  Perancangan", "Entity Relationship Diagram (ERD)")

    # Image area
    img_x = Inches(0.6)
    img_y = Inches(1.8)
    img_w = Inches(8.5)
    img_h = Inches(5.1)
    add_rect(s, img_x, img_y, img_w, img_h, fill=C_CARD_BG)
    try:
        # ERD image is wide (751x372) - fit to box
        s.shapes.add_picture(os.path.join(DIAGRAM_DIR, "ERD/erd.png"),
                              img_x + Inches(0.2), img_y + Inches(0.6),
                              width=img_w - Inches(0.4))
    except Exception as e:
        add_text(s, img_x, img_y, img_w, img_h, "ERD",
                 font_size=14, color=C_MUTED, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Right side: entity summary
    right_x = Inches(9.35)
    right_w = Inches(3.4)
    add_text(s, right_x, Inches(1.8), right_w, Inches(0.4),
             "ENTITAS UTAMA", font_size=11, color=C_ACCENT,
             font_name="Calibri", bold=True)
    add_rect(s, right_x, Inches(2.15), Inches(0.5), Inches(0.04),
             fill=C_ACCENT)

    entities = [
        ("Gejala", "30 data", C_MINT),
        ("Penyakit", "10 data", C_ACCENT),
        ("Aturan", "30 data", C_GOLD),
        ("Detail Aturan", "90 relasi", C_NAVY),
        ("Diagnosa", "Riwayat", C_PRIMARY),
        ("Admin", "1 akun", C_RED),
    ]
    y = Inches(2.4)
    for i, (name, info, color) in enumerate(entities):
        y_pos = y + Inches(0.7) * i
        add_rect(s, right_x, y_pos, right_w, Inches(0.6), fill=C_CARD_BG)
        add_rect(s, right_x, y_pos, Inches(0.08), Inches(0.6), fill=color)
        add_text(s, right_x + Inches(0.25), y_pos + Inches(0.07),
                 right_w - Inches(0.3), Inches(0.3), name,
                 font_size=12, color=C_PRIMARY, font_name="Calibri",
                 bold=True)
        add_text(s, right_x + Inches(0.25), y_pos + Inches(0.32),
                 right_w - Inches(0.3), Inches(0.3), info,
                 font_size=9, color=C_MUTED, font_name="Calibri",
                 italic=True)

    add_slide_number(s, 12)


# ============================================================
# SLIDE 13: DEMO
# ============================================================
def slide_13_demo():
    s = new_slide(bg=C_PRIMARY_DARK)

    # Top accent
    add_rect(s, 0, 0, W, Inches(0.15), fill=C_GOLD)

    # Header
    add_text(s, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.18),
             "", font_size=11)
    add_rect(s, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.18),
             fill=C_ACCENT)
    add_text(s, Inches(0.85), Inches(0.45), Inches(8), Inches(0.35),
             "DEMO SISTEM", font_size=11, color=C_ACCENT,
             font_name="Calibri", bold=True)
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.7),
             "Demonstrasi Aplikasi Klinik Merpati",
             font_size=32, color=C_LIGHT_BG, font_name="Georgia", bold=True)

    # Decorative
    add_oval(s, Inches(11.5), Inches(0.3), Inches(2.5), Inches(2.5),
             fill=C_PRIMARY, line=None)

    # Mock browser window
    bx = Inches(0.6)
    by = Inches(1.85)
    bw = Inches(8.5)
    bh = Inches(4.4)
    add_rect(s, bx, by, bw, bh, fill=C_LIGHT_BG)

    # Browser top bar
    add_rect(s, bx, by, bw, Inches(0.45), fill=C_PRIMARY)
    add_oval(s, bx + Inches(0.2), by + Inches(0.12), Inches(0.2),
             Inches(0.2), fill=C_ACCENT)
    add_oval(s, bx + Inches(0.5), by + Inches(0.12), Inches(0.2),
             Inches(0.2), fill=C_GOLD)
    add_oval(s, bx + Inches(0.8), by + Inches(0.12), Inches(0.2),
             Inches(0.2), fill=C_MINT)
    add_rect(s, bx + Inches(1.3), by + Inches(0.1), Inches(6.5),
             Inches(0.25), fill=C_PRIMARY_DARK)
    add_text(s, bx + Inches(1.4), by + Inches(0.1), Inches(6.4),
             Inches(0.25), "klinik-merpati.local/konsultasi.php",
             font_size=9, color=C_LIGHT_BG, font_name="Calibri",
             anchor=MSO_ANCHOR.MIDDLE, italic=True)

    # Browser body — mock consultation
    add_text(s, bx + Inches(0.4), by + Inches(0.7), bw - Inches(0.8),
             Inches(0.5), "Analisis Kesehatan Avian", font_size=22,
             color=C_PRIMARY, font_name="Georgia", bold=True)
    add_text(s, bx + Inches(0.4), by + Inches(1.2), bw - Inches(0.8),
             Inches(0.5),
             "Pilih gejala yang teramati pada merpati Anda:",
             font_size=12, color=C_MUTED, font_name="Calibri")

    # Mock name input
    add_text(s, bx + Inches(0.4), by + Inches(1.7), Inches(2.5),
             Inches(0.3), "Nama Pemilik", font_size=10, color=C_PRIMARY,
             font_name="Calibri", bold=True)
    add_rect(s, bx + Inches(0.4), by + Inches(2.0), Inches(7.7),
             Inches(0.4), fill=C_WHITE, line=C_LINE, line_width=0.75)
    add_text(s, bx + Inches(0.55), by + Inches(2.0), Inches(7.5),
             Inches(0.4), "Pak Hadi - Peternak Merpati",
             font_size=10, color=C_DARK, font_name="Calibri",
             anchor=MSO_ANCHOR.MIDDLE, italic=True)

    # Mock checkboxes (2 cols x 3 rows)
    cb_x = bx + Inches(0.4)
    cb_y = by + Inches(2.55)
    cb_w = Inches(3.7)
    cb_h = Inches(0.32)
    cb_gap_x = Inches(0.2)
    cb_gap_y = Inches(0.08)
    options = [
        ("G03", "Diare", True),
        ("G05", "Berat badan menurun", True),
        ("G04", "Diare berdarah", False),
        ("G25", "Kotoran berlendir", True),
        ("G01", "Nafsu makan menurun", True),
        ("G02", "Burung terlihat lesu", False),
    ]
    for i, (code, label, checked) in enumerate(options):
        col = i % 2
        row = i // 2
        x = cb_x + (cb_w + cb_gap_x) * col
        y = cb_y + (cb_h + cb_gap_y) * row
        add_rect(s, x, y, cb_w, cb_h, fill=C_CARD_BG)
        # Checkbox
        add_rect(s, x + Inches(0.1), y + Inches(0.07), Inches(0.18),
                 Inches(0.18), fill=(C_PRIMARY if checked else C_WHITE),
                 line=C_PRIMARY, line_width=0.75)
        if checked:
            add_text(s, x + Inches(0.1), y + Inches(0.03), Inches(0.18),
                     Inches(0.22), "✓", font_size=10, color=C_LIGHT_BG,
                     font_name="Calibri", bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.38), y, Inches(3.2), cb_h,
                 f"{code} - {label}", font_size=10, color=C_DARK,
                 font_name="Calibri", anchor=MSO_ANCHOR.MIDDLE)

    # Diagnosa button
    btn_x = bx + Inches(0.4)
    btn_y = by + Inches(3.78)
    add_rect(s, btn_x, btn_y, Inches(2.2), Inches(0.45), fill=C_ACCENT)
    add_text(s, btn_x, btn_y, Inches(2.2), Inches(0.45), "JALANKAN DIAGNOSA",
             font_size=10, color=C_LIGHT_BG, font_name="Calibri",
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Right panel: result preview
    rx = Inches(9.4)
    ry = Inches(1.85)
    rw = Inches(3.4)
    rh = Inches(4.4)
    add_rect(s, rx, ry, rw, rh, fill=C_PRIMARY)

    add_text(s, rx + Inches(0.3), ry + Inches(0.25), Inches(2),
             Inches(0.3), "HASIL DIAGNOSA", font_size=10, color=C_GOLD,
             font_name="Calibri", bold=True)
    add_rect(s, rx + Inches(0.3), ry + Inches(0.6), Inches(0.4),
             Inches(0.04), fill=C_GOLD)

    add_text(s, rx + Inches(0.3), ry + Inches(0.8), rw - Inches(0.6),
             Inches(0.3), "Tingkat Kecocokan", font_size=12, color=C_LIGHT_BG,
             font_name="Calibri", line_spacing=1.1)
    add_text(s, rx + Inches(0.3), ry + Inches(1.15), rw - Inches(0.6),
             Inches(0.9), "100%", font_size=48, color=C_GOLD,
             font_name="Georgia", bold=True, line_spacing=1.0)
    add_rect(s, rx + Inches(0.3), ry + Inches(2.1), rw - Inches(0.6),
             Inches(0.04), fill=C_LIGHT_BG)
    add_text(s, rx + Inches(0.3), ry + Inches(2.25), rw - Inches(0.6),
             Inches(0.3), "Penyakit Terdeteksi", font_size=10, color=C_GOLD,
             font_name="Calibri", bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(2.5), rw - Inches(0.6),
             Inches(0.45), "Coccidiosis", font_size=18, color=C_LIGHT_BG,
             font_name="Georgia", bold=True)
    add_text(s, rx + Inches(0.3), ry + Inches(3.0), rw - Inches(0.6),
             Inches(0.3), "(P03)", font_size=10, color=C_GOLD,
             font_name="Calibri", italic=True)
    add_text(s, rx + Inches(0.3), ry + Inches(3.4), rw - Inches(0.6),
             Inches(0.9), "Solusi & Pencegahan tersedia secara lengkap "
             "pada halaman hasil diagnosa.",
             font_size=9, color=C_SECONDARY, font_name="Calibri",
             line_spacing=1.35, italic=True)

    add_slide_number(s, 13)


# ============================================================
# SLIDE 14: TABEL PENGUJIAN
# ============================================================
def slide_14_tabel_pengujian():
    s = new_slide(bg=C_LIGHT_BG)
    add_header_band(s, "BAB V  •  Pengujian", "Tabel Pengujian Black-Box")

    # Stats summary
    stats = [
        ("19", "Skenario Pengujian", C_ACCENT),
        ("100%", "Berhasil", C_MINT),
        ("0", "Gagal", C_RED),
    ]
    sw = Inches(1.3)
    sh = Inches(0.75)
    sx = Inches(8.5)
    sy = Inches(1.7)
    for i, (num, label, color) in enumerate(stats):
        x = sx + (sw + Inches(0.1)) * i
        add_rect(s, x, sy, sw, sh, fill=color)
        add_text(s, x, sy + Inches(0.03), sw, Inches(0.4), num,
                 font_size=18, color=C_LIGHT_BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name="Georgia", bold=True)
        add_text(s, x, sy + Inches(0.42), sw, Inches(0.33), label,
                 font_size=8, color=C_LIGHT_BG, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font_name="Calibri", bold=True)

    # Table
    rows_data = [
        ("1", "Login Admin", "Berhasil"),
        ("2", "Login Gagal", "Berhasil"),
        ("3", "Logout Admin", "Berhasil"),
        ("4", "Tambah Gejala", "Berhasil"),
        ("5", "Edit Gejala", "Berhasil"),
        ("6", "Hapus Gejala", "Berhasil"),
        ("7", "Tambah Penyakit", "Berhasil"),
        ("8", "Konsultasi", "Berhasil"),
        ("9", "Hasil Diagnosa", "Berhasil"),
        ("10", "Cari Riwayat", "Berhasil"),
    ]
    tx = Inches(0.6)
    ty = Inches(2.6)
    tw = Inches(12.1)
    rh = Inches(0.32)
    no_w = Inches(0.7)
    name_w = Inches(8.4)
    res_w = Inches(3.0)

    # Header
    add_rect(s, tx, ty, tw, Inches(0.4), fill=C_PRIMARY)
    add_text(s, tx + Inches(0.2), ty, no_w, Inches(0.4), "NO",
             font_size=11, color=C_GOLD, font_name="Calibri", bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + no_w + Inches(0.2), ty, name_w, Inches(0.4),
             "PENGUJIAN", font_size=11, color=C_GOLD, font_name="Calibri",
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + no_w + name_w + Inches(0.2), ty, res_w, Inches(0.4),
             "HASIL", font_size=11, color=C_GOLD, font_name="Calibri",
             bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    for i, (no, name, result) in enumerate(rows_data):
        y_row = ty + Inches(0.4) + rh * i
        bg = C_WHITE if i % 2 == 0 else C_CARD_BG
        add_rect(s, tx, y_row, tw, rh, fill=bg)
        add_text(s, tx + Inches(0.2), y_row, no_w, rh, no,
                 font_size=10, color=C_PRIMARY, font_name="Calibri",
                 bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + no_w + Inches(0.2), y_row, name_w, rh, name,
                 font_size=10, color=C_DARK, font_name="Calibri",
                 anchor=MSO_ANCHOR.MIDDLE)
        # Result pill
        pill_x = tx + no_w + name_w + Inches(0.4)
        pill_w = Inches(2.2)
        pill_h = Inches(0.24)
        pill_y = y_row + (rh - pill_h) / 2
        add_rect(s, pill_x, pill_y, pill_w, pill_h, fill=C_MINT)
        add_text(s, pill_x, pill_y, pill_w, pill_h, "✓ " + result,
                 font_size=9, color=C_WHITE, font_name="Calibri",
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom callout
    note_y = ty + Inches(0.4) + rh * len(rows_data) + Inches(0.2)
    add_rect(s, Inches(0.6), note_y, Inches(12.1), Inches(0.45),
             fill=C_PRIMARY)
    add_text(s, Inches(0.85), note_y, Inches(11.6), Inches(0.45),
             "Seluruh 19 skenario pengujian Black-Box menunjukkan "
             "status BERHASIL — sistem berfungsi sesuai spesifikasi.",
             font_size=11, color=C_LIGHT_BG, font_name="Calibri",
             italic=True, anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, 14)


# ============================================================
# SLIDE 15: KESIMPULAN
# ============================================================
def slide_15_kesimpulan():
    s = new_slide(bg=C_PRIMARY_DARK)

    # Top accent
    add_rect(s, 0, 0, W, Inches(0.15), fill=C_GOLD)

    # Decorative
    add_oval(s, Inches(10.5), Inches(-2.5), Inches(5), Inches(5),
             fill=C_PRIMARY, line=None)
    add_oval(s, Inches(-1.5), Inches(5.5), Inches(3.5), Inches(3.5),
             fill=C_PRIMARY, line=None)

    # Header
    add_rect(s, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.18),
             fill=C_ACCENT)
    add_text(s, Inches(0.85), Inches(0.45), Inches(8), Inches(0.35),
             "BAB VI  •  Penutup", font_size=11, color=C_ACCENT,
             font_name="Calibri", bold=True)
    add_text(s, Inches(0.6), Inches(0.85), Inches(11), Inches(0.7),
             "Kesimpulan & Saran",
             font_size=32, color=C_LIGHT_BG, font_name="Georgia", bold=True)

    # Two columns
    col_w = Inches(5.9)
    col_h = Inches(5.0)
    y = Inches(1.85)
    gap = Inches(0.3)
    x1 = Inches(0.6)
    x2 = x1 + col_w + gap

    # === KESIMPULAN ===
    add_rect(s, x1, y, col_w, col_h, fill=C_CARD_BG)
    add_rect(s, x1, y, col_w, Inches(0.6), fill=C_GOLD)
    add_text(s, x1 + Inches(0.3), y, Inches(2), Inches(0.6), "01",
             font_size=20, color=C_PRIMARY_DARK, font_name="Georgia",
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x1, y, col_w, Inches(0.6), "KESIMPULAN",
             font_size=18, color=C_PRIMARY_DARK, font_name="Georgia",
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    conclusions = [
        "Sistem pakar Klinik Merpati berhasil dibangun berbasis web "
        "dengan PHP dan MySQL.",
        "Metode Forward Chaining terbukti efektif dalam mencocokkan "
        "gejala dengan basis aturan untuk menghasilkan diagnosa.",
        "Seluruh 19 skenario pengujian Black-Box menunjukkan hasil "
        "BERHASIL dengan tingkat validitas yang baik.",
        "Sistem mampu membantu peternak merpati melakukan deteksi "
        "dini secara mandiri, cepat, dan konsisten.",
    ]
    cy = y + Inches(0.85)
    for i, c in enumerate(conclusions):
        # Bullet
        add_oval(s, x1 + Inches(0.3), cy + Inches(0.15), Inches(0.15),
                 Inches(0.15), fill=C_ACCENT)
        add_text(s, x1 + Inches(0.6), cy, col_w - Inches(0.8),
                 Inches(0.95), c, font_size=11, color=C_DARK,
                 font_name="Calibri", line_spacing=1.4)
        cy += Inches(0.95)

    # === SARAN ===
    add_rect(s, x2, y, col_w, col_h, fill=C_CARD_BG)
    add_rect(s, x2, y, col_w, Inches(0.6), fill=C_ACCENT)
    add_text(s, x2 + Inches(0.3), y, Inches(2), Inches(0.6), "02",
             font_size=20, color=C_LIGHT_BG, font_name="Georgia",
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x2, y, col_w, Inches(0.6), "SARAN PENGEMBANGAN",
             font_size=18, color=C_LIGHT_BG, font_name="Georgia",
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    suggestions = [
        "Penambahan jumlah penyakit dan gejala pada basis pengetahuan "
        "untuk cakupan yang lebih luas.",
        "Implementasi metode Certainty Factor untuk menangani "
        "ketidakpastian dalam diagnosa.",
        "Pengembangan aplikasi mobile berbasis Android atau iOS "
        "agar lebih mudah diakses di lapangan.",
        "Integrasi dengan basis data veteriner nasional dan fitur "
        "tele-konsultasi dengan dokter hewan.",
    ]
    cy = y + Inches(0.85)
    for i, sug in enumerate(suggestions):
        add_oval(s, x2 + Inches(0.3), cy + Inches(0.15), Inches(0.15),
                 Inches(0.15), fill=C_GOLD)
        add_text(s, x2 + Inches(0.6), cy, col_w - Inches(0.8),
                 Inches(0.95), sug, font_size=11, color=C_DARK,
                 font_name="Calibri", line_spacing=1.4)
        cy += Inches(0.95)

    # Bottom thank you bar
    add_rect(s, 0, H - Inches(0.7), W, Inches(0.7), fill=C_GOLD)
    add_text(s, Inches(0.6), H - Inches(0.7), Inches(8), Inches(0.7),
             "TERIMA KASIH  •  Thank You", font_size=18, color=C_PRIMARY_DARK,
             font_name="Georgia", bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.5), H - Inches(0.7), Inches(4.3), Inches(0.7),
             "Klinik Merpati  •  2025", font_size=11, color=C_PRIMARY_DARK,
             font_name="Calibri", bold=True, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.RIGHT)


# ============================================================
# GENERATE
# ============================================================
slide_01_cover()
slide_02_latar_belakang()
slide_03_rumusan_masalah()
slide_04_batasan_masalah()
slide_05_tujuan()
slide_06_manfaat()
slide_07_metode_data()
slide_08_metode_pengembangan()
slide_09_perancangan()
slide_10_use_case()
slide_11_activity()
slide_12_erd()
slide_13_demo()
slide_14_tabel_pengujian()
slide_15_kesimpulan()

prs.save(OUTPUT)
print(f"[OK] Saved: {OUTPUT}")
print(f"  Total slides: {len(prs.slides)}")
