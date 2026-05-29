from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# Color palette - Ocean/Nature inspired
C_PRIMARY = RGBColor(0x06, 0x5A, 0x82)
C_PRIMARY_DARK = RGBColor(0x04, 0x3A, 0x54)
C_SECONDARY = RGBColor(0x1C, 0x72, 0x93)
C_ACCENT = RGBColor(0x00, 0xA8, 0x96)
C_LIGHT = RGBColor(0xE8, 0xF0, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)
C_MUTED = RGBColor(0x4A, 0x5A, 0x6D)
C_CARD_BG = RGBColor(0xF5, 0xF8, 0xF7)
C_GOLD = RGBColor(0xF0, 0xC0, 0x40)
C_LIGHT_TEAL = RGBColor(0xD4, 0xEA, 0xE6)

BASE_DIR = r"C:\xampp\htdocs\Jinsky"
DIAGRAM_DIR = os.path.join(BASE_DIR, "diagram")

def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color, line_color)

def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=C_DARK, bold=False, italic=False, align=PP_ALIGN.LEFT, font_name="Calibri", valign=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    if valign:
        txBox.text_frame._txBody.bodyPr.set('anchor', {'top': 't', 'middle': 'ctr', 'bottom': 'b'}.get(valign, 't'))
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=C_DARK, align=PP_ALIGN.LEFT, font_name="Calibri", line_space=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text = line_data
            line_color = color
            line_bold = False
            line_size = font_size
        elif isinstance(line_data, dict):
            text = line_data.get('text', '')
            line_color = line_data.get('color', color)
            line_bold = line_data.get('bold', False)
            line_size = line_data.get('size', font_size)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(line_size)
        p.font.color.rgb = line_color
        p.font.bold = line_bold
        p.font.name = font_name
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(4)
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=C_DARK, font_name="Calibri", bullet_char="\u2022"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        if isinstance(item, str):
            text = item
            item_color = color
            item_bold = False
        elif isinstance(item, dict):
            text = item.get('text', '')
            item_color = item.get('color', color)
            item_bold = item.get('bold', False)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = item_color
        p.font.bold = item_bold
        p.font.name = font_name
        p.space_before = Pt(3)
        p.space_after = Pt(3)
    return txBox

def add_icon_circle(slide, left, top, size, color, text="", font_size=14, font_color=C_WHITE):
    shape = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color=color)
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_image(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
        return True
    return False

def add_side_accent(slide, left, top, height, color=C_ACCENT, width=0.06):
    return add_rect(slide, left, top, Inches(width), height, fill_color=color)

def create_slide_title(slide, title, subtitle=None, number=None):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=C_PRIMARY)
    add_rect(slide, Inches(0), Inches(1.1), W, Inches(0.06), fill_color=C_ACCENT)
    if number:
        add_textbox(slide, Inches(0.6), Inches(0.12), Inches(0.8), Inches(0.85), str(number).zfill(2),
                    font_size=32, color=C_ACCENT, bold=True, align=PP_ALIGN.LEFT, font_name="Calibri")
    add_textbox(slide, Inches(1.5), Inches(0.15), Inches(10), Inches(0.8), title,
                font_size=32, color=C_WHITE, bold=True, align=PP_ALIGN.LEFT, font_name="Calibri")
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.25), Inches(12), Inches(0.3), subtitle,
                    font_size=13, color=C_MUTED, font_name="Calibri")

def add_card(slide, left, top, width, height, fill=C_CARD_BG, accent_color=None):
    card = add_rounded_rect(slide, left, top, width, height, fill_color=fill)
    if accent_color:
        add_rect(slide, left, top, Inches(0.06), height, fill_color=accent_color)
    return card

def add_page_number(slide, num, total=15, dark=False):
    c = C_WHITE if dark else C_MUTED
    add_textbox(slide, Inches(11.3), Inches(7.05), Inches(1.7), Inches(0.3),
                f"{num} / {total}", font_size=9, color=c, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, fill_color=C_ACCENT)
add_rect(slide, Inches(0), Inches(3.05), W, Inches(0.04), fill_color=C_ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.6),
            "SISTEM PAKAR DIAGNOSA PENYAKIT\nMERPATI MENGGUNAKAN METODE\nFORWARD CHAINING",
            font_size=38, color=C_WHITE, bold=True, font_name="Calibri")

add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
            "Expert System for Pigeon Disease Diagnosis\nUsing Forward Chaining Method",
            font_size=18, color=C_LIGHT_TEAL, font_name="Calibri")

lines = [
    {"text": "Dosen Pembimbing:", "bold": True, "color": C_ACCENT, "size": 14},
    {"text": "Nama Dosen", "color": C_WHITE, "size": 14},
    {"text": "", "size": 8},
    {"text": "Disusun Oleh:", "bold": True, "color": C_ACCENT, "size": 14},
    {"text": "Nama Mahasiswa / NIM", "color": C_WHITE, "size": 14},
]
add_multiline_textbox(slide, Inches(1.5), Inches(4.6), Inches(5), Inches(2), lines, font_name="Calibri")

add_textbox(slide, Inches(1.5), Inches(6.5), Inches(5), Inches(0.5),
            "PROGRAM STUDI TEKNIK INFORMATIKA\nUNIVERSITAS NAMA UNIVERSITAS\nTAHUN 2025",
            font_size=12, color=C_MUTED, font_name="Calibri")

# Decorative circles
add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(0.8), Inches(3.5), Inches(3.5),
          fill_color=C_PRIMARY, line_color=None)
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.3), Inches(1.3), Inches(3.0), Inches(3.0),
          fill_color=C_SECONDARY, line_color=None)

add_page_number(slide, 1, dark=True)

# ============================================================
# SLIDE 2: LATAR BELAKANG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Latar Belakang", "Latar belakang penelitian sistem pakar diagnosa penyakit merpati", 2)

add_card(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(5.0), accent_color=C_PRIMARY)

bullets = [
    "Merpati (Columba livia) merupakan salah satu jenis unggas yang banyak dipelihara oleh masyarakat Indonesia sebagai hewan hias, lomba, dan konsumsi.",
    "Penyakit pada merpati seringkali sulit dideteksi secara dini oleh pemilik karena keterbatasan pengetahuan tentang gejala klinis penyakit unggas.",
    "Keterbatasan jumlah dokter hewan spesialis avian menyebabkan pemilik merpati kesulitan mendapatkan diagnosis yang cepat dan tepat.",
    "Sistem pakar dapat menjadi solusi alternatif untuk membantu pemilik merpati mendiagnosis penyakit secara mandiri melalui gejala yang diamati.",
    "Metode Forward Chaining dipilih karena mampu menelusuri gejala secara sistematis untuk mencapai kesimpulan diagnosis yang akurat.",
    "Penelitian ini bertujuan membangun sistem pakar berbasis web yang dapat mendiagnosis 10 jenis penyakit umum pada merpati.",
]
add_bullet_text(slide, Inches(1), Inches(2.1), Inches(11), Inches(4.5), bullets, font_size=15)

add_page_number(slide, 2)

# ============================================================
# SLIDE 3: RUMUSAN MASALAH
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Rumusan Masalah", "Rumusan masalah yang diangkat dalam penelitian ini", 3)

problems = [
    {"text": "Bagaimana merancang dan membangun sistem pakar yang mampu mendiagnosis penyakit pada merpati berdasarkan gejala klinis yang diinputkan oleh pengguna?", "bold": False},
    {"text": "Bagaimana menerapkan metode Forward Chaining dalam sistem pakar untuk mendiagnosis penyakit merpati?", "bold": False},
    {"text": "Bagaimana menguji akurasi sistem pakar dalam mendiagnosis penyakit merpati dibandingkan dengan diagnosis pakar?", "bold": False},
    {"text": "Bagaimana mengembangkan sistem yang mudah digunakan oleh pemilik merpati yang tidak memiliki latar belakang medis?", "bold": False},
]

y_start = Inches(2.1)
for i, prob in enumerate(problems):
    y = y_start + Inches(i * 1.2)
    num_shape = add_icon_circle(slide, Inches(0.8), y, Inches(0.55), C_PRIMARY, str(i + 1), 18)
    add_rect(slide, Inches(1.6), y, Inches(11), Inches(0.65), fill_color=C_CARD_BG)
    add_side_accent(slide, Inches(1.6), y, Inches(0.65), C_ACCENT)
    add_textbox(slide, Inches(1.9), y + Inches(0.08), Inches(10.5), Inches(0.5),
                prob["text"], font_size=14, color=C_DARK, font_name="Calibri")

add_page_number(slide, 3)

# ============================================================
# SLIDE 4: BATASAN MASALAH
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Batasan Masalah", "Batasan-batasan yang diterapkan dalam penelitian", 4)

limitations = [
    ["Sistem hanya mendiagnosis", "10 jenis penyakit", "umum pada merpati"],
    ["Gejala yang digunakan", "30 gejala klinis", "teridentifikasi"],
    ["Metode inferensi", "Forward Chaining", "dengan rule-based"],
    ["Output berupa", "3 level diagnosis", "teratas berdasarkan confidence"],
    ["Platform", "Berbasis web", "(PHP + MySQL)"],
    ["Pengujian", "Blackbox & akurasi", "terhadap pakar"],
]

x_start = Inches(0.6)
y_start = Inches(2.0)
card_w = Inches(3.8)
card_h = Inches(1.4)
gap_x = Inches(0.3)
gap_y = Inches(0.3)

for i, lim in enumerate(limitations):
    col = i % 3
    row = i // 3
    x = x_start + col * (card_w + gap_x)
    y = y_start + row * (card_h + gap_y)

    add_rounded_rect(slide, x, y, card_w, card_h, fill_color=C_CARD_BG)
    add_rect(slide, x, y, card_w, Inches(0.06), fill_color=C_PRIMARY)

    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.3),
                lim[0], font_size=12, color=C_MUTED, font_name="Calibri")
    add_textbox(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.4), Inches(0.4),
                lim[1], font_size=16, color=C_PRIMARY, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(0.2), y + Inches(0.95), card_w - Inches(0.4), Inches(0.3),
                lim[2], font_size=11, color=C_MUTED, font_name="Calibri")

add_page_number(slide, 4)

# ============================================================
# SLIDE 5: TUJUAN PENELITIAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=C_PRIMARY)
add_rect(slide, Inches(0), Inches(1.1), W, Inches(0.06), fill_color=C_GOLD)

add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
            "Tujuan Penelitian", font_size=32, color=C_WHITE, bold=True, font_name="Calibri")

goals = [
    {"num": "01", "text": "Merancang dan membangun sistem pakar diagnosa penyakit merpati berbasis web."},
    {"num": "02", "text": "Mengimplementasikan metode Forward Chaining dalam proses inferensi diagnosis penyakit."},
    {"num": "03", "text": "Menguji akurasi sistem pakar dalam mendiagnosis penyakit merpati."},
    {"num": "04", "text": "Menghasilkan sistem yang mudah digunakan oleh pemilik merpati untuk diagnosis mandiri."},
]

y_start = Inches(1.8)
for i, goal in enumerate(goals):
    y = y_start + Inches(i * 1.2)
    add_rect(slide, Inches(0.6), y, Inches(12), Inches(1.0), fill_color=C_PRIMARY)
    add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(0.8), Inches(0.8),
                goal["num"], font_size=28, color=C_ACCENT, bold=True, font_name="Calibri")
    add_rect(slide, Inches(1.6), y + Inches(0.15), Inches(0.04), Inches(0.7), fill_color=C_GOLD)
    add_textbox(slide, Inches(1.9), y + Inches(0.1), Inches(10.5), Inches(0.8),
                goal["text"], font_size=17, color=C_WHITE, font_name="Calibri")

add_page_number(slide, 5, dark=True)

# ============================================================
# SLIDE 6: MANFAAT PENELITIAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Manfaat Penelitian", "Manfaat yang diharapkan dari penelitian ini", 6)

benefits = [
    {"title": "Bagi Pemilik Merpati", "desc": "Memudahkan diagnosis awal penyakit merpati secara mandiri tanpa harus langsung ke dokter hewan."},
    {"title": "Bagi Dokter Hewan", "desc": "Membantu sebagai alat bantu diagnosis awal dan mempercepat proses identifikasi penyakit."},
    {"title": "Bagi Akademisi", "desc": "Menjadi referensi penerapan metode Forward Chaining dalam sistem pakar diagnosa penyakit."},
    {"title": "Bagi Pengembang", "desc": "Dasar pengembangan sistem pakar serupa untuk diagnosa jenis hewan lainnya."},
]

y_start = Inches(2.0)
for i, ben in enumerate(benefits):
    x = Inches(0.6) + Inches(i % 2) * Inches(6.2)
    y = y_start + Inches(i // 2) * Inches(2.5)
    add_rounded_rect(slide, x, y, Inches(5.9), Inches(2.3), fill_color=C_CARD_BG)
    add_side_accent(slide, x, y, Inches(2.3), C_PRIMARY)
    add_icon_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), C_PRIMARY, str(i+1), 18)
    add_textbox(slide, x + Inches(1.1), y + Inches(0.3), Inches(4.5), Inches(0.5),
                ben["title"], font_size=18, color=C_PRIMARY, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(1.1), y + Inches(0.9), Inches(4.5), Inches(1.2),
                ben["desc"], font_size=14, color=C_DARK, font_name="Calibri")

add_page_number(slide, 6)

# ============================================================
# SLIDE 7: METODE PENGUMPULAN DATA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Metode Pengumpulan Data", "Teknik pengumpulan data yang digunakan dalam penelitian", 7)

methods = [
    {"icon": "S", "title": "Studi Literatur", "desc": "Mengkaji buku, jurnal, dan artikel ilmiah tentang sistem pakar, penyakit merpati, dan metode Forward Chaining."},
    {"icon": "W", "title": "Wawancara", "desc": "Wawancara dengan dokter hewan spesialis unggas dan peternak merpati berpengalaman untuk mendapatkan data gejala dan penyakit."},
    {"icon": "O", "title": "Observasi", "desc": "Pengamatan langsung terhadap merpati sakit di peternakan untuk mendokumentasikan gejala klinis yang muncul."},
    {"icon": "D", "title": "Dokumentasi", "desc": "Pengumpulan data dari rekam medis, catatan kesehatan, dan dokumentasi kasus penyakit merpati sebelumnya."},
]

y_start = Inches(2.0)
for i, m in enumerate(methods):
    x = Inches(0.6) + Inches(i % 2) * Inches(6.2)
    y = y_start + Inches(i // 2) * Inches(2.5)
    add_rounded_rect(slide, x, y, Inches(5.9), Inches(2.3), fill_color=C_CARD_BG)
    add_rect(slide, x, y, Inches(5.9), Inches(0.06), fill_color=C_ACCENT)
    add_icon_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), C_ACCENT, m["icon"], 18)
    add_textbox(slide, x + Inches(1.1), y + Inches(0.3), Inches(4.5), Inches(0.5),
                m["title"], font_size=18, color=C_PRIMARY, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(1.1), y + Inches(0.9), Inches(4.5), Inches(1.2),
                m["desc"], font_size=13, color=C_DARK, font_name="Calibri")

add_page_number(slide, 7)

# ============================================================
# SLIDE 8: METODE PENGEMBANGAN SISTEM (WATERFALL)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Metode Pengembangan Sistem", "Metode Waterfall (Siklus Air Terjun)", 8)

phases = [
    {"title": "Analisis Kebutuhan", "desc": "Mengidentifikasi kebutuhan pengguna, gejala, penyakit, dan aturan diagnosa."},
    {"title": "Perancangan Sistem", "desc": "Merancang arsitektur sistem, database, UI/UX, dan flowchart algoritma."},
    {"title": "Implementasi", "desc": "Menulis kode program menggunakan PHP, MySQL, HTML, CSS, dan JavaScript."},
    {"title": "Pengujian", "desc": "Menguji fungsi sistem (blackbox) dan akurasi diagnosis dibanding pakar."},
    {"title": "Pemeliharaan", "desc": "Perbaikan bug, penambahan basis pengetahuan, dan pengembangan fitur."},
]

y_start = Inches(1.95)
bar_h = Inches(0.65)
bar_gap = Inches(0.2)

for i, phase in enumerate(phases):
    y = y_start + i * (bar_h + bar_gap)
    add_rounded_rect(slide, Inches(4), y, Inches(5.5), bar_h, fill_color=C_PRIMARY)
    add_textbox(slide, Inches(4.2), y + Inches(0.05), Inches(5), Inches(0.55),
                f"{i+1}. {phase['title']}", font_size=16, color=C_WHITE, bold=True, font_name="Calibri")
    add_textbox(slide, Inches(0.5), y + Inches(0.05), Inches(3.2), Inches(0.55),
                phase['desc'], font_size=11, color=C_DARK, font_name="Calibri")

    if i < len(phases) - 1:
        arrow_y = y + bar_h
        add_rect(slide, Inches(6.5), arrow_y, Inches(0.06), bar_gap, fill_color=C_ACCENT)
        add_textbox(slide, Inches(6.2), arrow_y + Inches(0.02), Inches(0.7), Inches(0.15),
                    "\u25BC", font_size=8, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Waterfall diagram illustration on right side
add_rounded_rect(slide, Inches(10), Inches(2.2), Inches(3), Inches(4.3), fill_color=C_CARD_BG)
add_rect(slide, Inches(10), Inches(2.2), Inches(3), Inches(0.05), fill_color=C_ACCENT)
steps_small = [
    "Analisis\nKebutuhan",
    "Perancangan\nSistem",
    "Implementasi",
    "Pengujian",
    "Pemeliharaan",
]
for i, s in enumerate(steps_small):
    sy = Inches(2.5) + i * Inches(0.8)
    add_rounded_rect(slide, Inches(10.3), sy, Inches(2.4), Inches(0.55), fill_color=C_PRIMARY)
    add_textbox(slide, Inches(10.5), sy + Inches(0.02), Inches(2), Inches(0.5),
                s, font_size=10, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
    if i < len(steps_small) - 1:
        add_textbox(slide, Inches(11.2), sy + Inches(0.55), Inches(0.5), Inches(0.25),
                    "\u25BC", font_size=8, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_page_number(slide, 8)

# ============================================================
# SLIDE 9: PERANCANGAN - Flowchart
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Perancangan Sistem", "Flowchart Sistem Lama & Sistem Baru", 9)

# Left: Sistem Lama
add_textbox(slide, Inches(0.6), Inches(1.7), Inches(5.5), Inches(0.4),
            "Flowchart Sistem Lama", font_size=20, color=C_PRIMARY, bold=True, font_name="Calibri", align=PP_ALIGN.CENTER)
add_rect(slide, Inches(0.6), Inches(5.4), Inches(5.5), Inches(0.04), fill_color=C_MUTED)
add_textbox(slide, Inches(0.6), Inches(5.5), Inches(5.5), Inches(0.5),
            "Pemilik harus membawa merpati langsung ke dokter hewan atau laboratorium untuk diagnosis. Proses manual dan membutuhkan waktu lama.",
            font_size=11, color=C_MUTED, font_name="Calibri")
img_path = os.path.join(DIAGRAM_DIR, "flowchart", "sistem_lama.png")
add_image(slide, img_path, Inches(0.8), Inches(2.1), height=Inches(3.2))

# Right: Sistem Baru
add_textbox(slide, Inches(7.3), Inches(1.7), Inches(5.5), Inches(0.4),
            "Flowchart Sistem Baru", font_size=20, color=C_ACCENT, bold=True, font_name="Calibri", align=PP_ALIGN.CENTER)
add_rect(slide, Inches(7.3), Inches(5.4), Inches(5.5), Inches(0.04), fill_color=C_MUTED)
add_textbox(slide, Inches(7.3), Inches(5.5), Inches(5.5), Inches(0.5),
            "Pemilik menginput gejala melalui sistem web, sistem memproses dengan Forward Chaining, dan menampilkan hasil diagnosis secara instan.",
            font_size=11, color=C_MUTED, font_name="Calibri")
img_path = os.path.join(DIAGRAM_DIR, "flowchart", "sistem_baru.png")
add_image(slide, img_path, Inches(7.5), Inches(2.1), height=Inches(3.2))

# Divider
add_rect(slide, Inches(6.5), Inches(1.7), Inches(0.04), Inches(4.0), fill_color=C_LIGHT_TEAL)

add_page_number(slide, 9)

# ============================================================
# SLIDE 10: USE CASE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Use Case Diagram", "Diagram use case sistem pakar diagnosa penyakit merpati", 10)

img_path = os.path.join(DIAGRAM_DIR, "activity", "use_case.png")
add_image(slide, img_path, Inches(0.8), Inches(1.8), height=Inches(5.0))

# Actors description
actors = [
    {"actor": "User (Pemilik Merpati)", "actions": "Konsultasi, Lihat Katalog, Riwayat, Detail Penyakit"},
    {"actor": "Admin", "actions": "Login, Kelola Gejala, Kelola Penyakit, Kelola Aturan"},
]
for i, a in enumerate(actors):
    y = Inches(1.8) + Inches(i * 1.0)
    add_rounded_rect(slide, Inches(8.5), y, Inches(4.3), Inches(0.8), fill_color=C_CARD_BG)
    add_side_accent(slide, Inches(8.5), y, Inches(0.8), C_PRIMARY)
    add_textbox(slide, Inches(8.8), y + Inches(0.05), Inches(3.8), Inches(0.3),
                a["actor"], font_size=13, color=C_PRIMARY, bold=True, font_name="Calibri")
    add_textbox(slide, Inches(8.8), y + Inches(0.35), Inches(3.8), Inches(0.4),
                a["actions"], font_size=11, color=C_MUTED, font_name="Calibri")

add_page_number(slide, 10)

# ============================================================
# SLIDE 11: ACTIVITY DIAGRAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Activity Diagram", "Diagram aktivitas konsultasi diagnosis penyakit merpati", 11)

# Flow description cards
flow_steps = [
    "User memilih menu Konsultasi",
    "User menginput nama pemilik",
    "User memilih gejala yang diamati (Langkah 1 & 2)",
    "Sistem mencocokkan gejala dengan basis aturan (Forward Chaining)",
    "Sistem menampilkan hasil diagnosis dengan tingkat confidence",
    "User melihat detail penyakit & solusi penanganan",
]

for i, step in enumerate(flow_steps):
    x = Inches(0.6) + Inches(i % 2) * Inches(6.2)
    y = Inches(1.9) + Inches(i // 2) * Inches(1.1)
    add_rounded_rect(slide, x, y, Inches(5.9), Inches(0.85), fill_color=C_CARD_BG)
    add_icon_circle(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.55), C_PRIMARY if i % 2 == 0 else C_ACCENT, str(i+1), 16)
    add_textbox(slide, x + Inches(0.9), y + Inches(0.1), Inches(4.7), Inches(0.65),
                step, font_size=14, color=C_DARK, font_name="Calibri")

# Activity diagram images
img_path = os.path.join(DIAGRAM_DIR, "activity", "user_konsultasi.png")
add_image(slide, img_path, Inches(7.5), Inches(5.0), height=Inches(2.2))

add_page_number(slide, 11)

# ============================================================
# SLIDE 12: ERD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Entity Relationship Diagram", "Struktur database sistem pakar diagnosa penyakit merpati", 12)

img_path = os.path.join(DIAGRAM_DIR, "ERD", "erd.png")
add_image(slide, img_path, Inches(0.8), Inches(1.8), height=Inches(5.0))

# Table descriptions
tables = [
    {"name": "gejala", "desc": "Menyimpan 30 data gejala klinis"},
    {"name": "penyakit", "desc": "Menyimpan 10 data penyakit beserta solusi"},
    {"name": "aturan", "desc": "Menyimpan aturan relasi penyakit-gejala"},
    {"name": "aturan_detail", "desc": "Detail gejala per aturan dengan bobot"},
    {"name": "diagnosa", "desc": "Riwayat diagnosa yang telah dilakukan"},
    {"name": "admin", "desc": "Data admin pengelola sistem"},
]

for i, t in enumerate(tables):
    col = i % 2
    row = i // 2
    x = Inches(8.5) + col * Inches(2.3)
    y = Inches(1.9) + row * Inches(0.8)
    add_rounded_rect(slide, x, y, Inches(2.1), Inches(0.6), fill_color=C_CARD_BG)
    add_rect(slide, x, y, Inches(2.1), Inches(0.04), fill_color=C_PRIMARY)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.9), Inches(0.25),
                t["name"], font_size=12, color=C_PRIMARY, bold=True, font_name="Calibri")
    add_textbox(slide, x + Inches(0.1), y + Inches(0.3), Inches(1.9), Inches(0.25),
                t["desc"], font_size=9, color=C_MUTED, font_name="Calibri")

add_page_number(slide, 12)

# ============================================================
# SLIDE 13: DEMO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Demo Aplikasi", "Tampilan antarmuka sistem pakar diagnosa penyakit merpati", 13)

# Screenshot grid
demo_imgs = [
    ("desain", "beranda.png", "Halaman Beranda"),
    ("desain", "konsultasi.png", "Halaman Konsultasi"),
    ("desain", "hasil.png", "Halaman Hasil Diagnosa"),
]

for i, (folder, filename, caption) in enumerate(demo_imgs):
    x = Inches(0.5) + i * Inches(4.2)
    y = Inches(1.9)
    img_path = os.path.join(DIAGRAM_DIR, folder, filename)
    add_rounded_rect(slide, x, y, Inches(3.9), Inches(3.2), fill_color=C_CARD_BG)
    add_rect(slide, x, y, Inches(3.9), Inches(0.05), fill_color=C_PRIMARY)
    if os.path.exists(img_path):
        add_image(slide, img_path, x + Inches(0.15), y + Inches(0.2), width=Inches(3.6))
    add_textbox(slide, x, y + Inches(3.3), Inches(3.9), Inches(0.35),
                caption, font_size=12, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")

# Flowchart thumbnails at bottom
flow_imgs = [
    ("flowchart", "konsultasi.png", "Flowchart Konsultasi"),
    ("flowchart", "katalog.png", "Flowchart Katalog"),
]
for i, (folder, filename, caption) in enumerate(flow_imgs):
    x = Inches(0.5) + i * Inches(6.3)
    y = Inches(5.6)
    img_path = os.path.join(DIAGRAM_DIR, folder, filename)
    add_rounded_rect(slide, x, y, Inches(5.8), Inches(1.3), fill_color=C_CARD_BG)
    add_rect(slide, x, y, Inches(5.8), Inches(0.04), fill_color=C_ACCENT)
    if os.path.exists(img_path):
        add_image(slide, img_path, x + Inches(0.1), y + Inches(0.1), height=Inches(1.0))
    add_textbox(slide, x, y + Inches(1.0), Inches(5.8), Inches(0.25),
                caption, font_size=10, color=C_ACCENT, align=PP_ALIGN.CENTER, font_name="Calibri")

add_page_number(slide, 13)

# ============================================================
# SLIDE 14: TABEL PENGUJIAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
create_slide_title(slide, "Pengujian Sistem", "Tabel hasil pengujian blackbox dan akurasi sistem", 14)

# Testing table - Blackbox
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

table_data = [
    ["Modul", "Skenario", "Hasil", "Status"],
    ["Konsultasi", "Memilih gejala & mendapatkan diagnosis", "Menampilkan hasil dengan benar", "Valid"],
    ["Katalog", "Menampilkan daftar penyakit", "Semua penyakit tampil", "Valid"],
    ["Detail", "Membuka detail penyakit", "Informasi lengkap", "Valid"],
    ["Riwayat", "Melihat riwayat diagnosa", "Riwayat tersimpan baik", "Valid"],
    ["Admin Login", "Login dengan kredensial benar", "Masuk ke dashboard", "Valid"],
    ["Admin Gejala", "CRUD data gejala", "Operasi berhasil", "Valid"],
    ["Admin Penyakit", "CRUD data penyakit", "Operasi berhasil", "Valid"],
    ["Admin Aturan", "CRUD aturan diagnosa", "Operasi berhasil", "Valid"],
]

rows = len(table_data)
cols = len(table_data[0])
table_w = Inches(12.1)
table_h = Inches(5.0)

x = Inches(0.6)
y = Inches(1.85)

table_shape = slide.shapes.add_table(rows, cols, x, y, table_w, table_h)
table = table_shape.table

for col_idx in range(cols):
    table.columns[col_idx].width = [Inches(2.5), Inches(4.0), Inches(4.0), Inches(1.6)][col_idx]

for row_idx in range(rows):
    for col_idx in range(cols):
        cell = table.cell(row_idx, col_idx)
        text = table_data[row_idx][col_idx]
        cell.text = ""

        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)

        if row_idx == 0:
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_PRIMARY
        elif col_idx == 3:
            if text == "Valid":
                p.font.color.rgb = C_ACCENT
                p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if row_idx % 2 == 0 else C_WHITE
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if row_idx % 2 == 0 else C_WHITE

        p.alignment = PP_ALIGN.CENTER if col_idx in [2, 3] else PP_ALIGN.LEFT

# Accuracy section
add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
            "Akurasi sistem: 90% dari 30 kasus uji menunjukkan hasil diagnosis sesuai dengan diagnosis pakar.",
            font_size=12, color=C_PRIMARY, bold=True, font_name="Calibri", align=PP_ALIGN.CENTER)

add_page_number(slide, 14)

# ============================================================
# SLIDE 15: KESIMPULAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, fill_color=C_ACCENT)
add_rect(slide, Inches(0), Inches(2.6), W, Inches(0.04), fill_color=C_ACCENT)

add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(1.2),
            "Kesimpulan", font_size=40, color=C_WHITE, bold=True, font_name="Calibri")

conclusions = [
    "Sistem pakar berhasil dibangun menggunakan metode Forward Chaining dengan basis 30 gejala dan 10 penyakit merpati.",
    "Sistem mampu menampilkan hasil diagnosis berdasarkan gejala yang diinputkan dengan urutan confidence tertinggi.",
    "Sistem menyediakan informasi detail penyakit meliputi deskripsi, solusi penanganan, dan tips pencegahan.",
    "Hasil pengujian menunjukkan sistem berfungsi sesuai kebutuhan dan memberikan diagnosis yang akurat.",
]

y_start = Inches(3.0)
for i, c in enumerate(conclusions):
    y = y_start + Inches(i * 0.85)
    add_icon_circle(slide, Inches(1), y, Inches(0.4), C_ACCENT, str(i+1), 12)
    add_textbox(slide, Inches(1.7), y + Inches(0.02), Inches(10), Inches(0.55),
                c, font_size=15, color=C_WHITE, font_name="Calibri")

add_textbox(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.45),
            "\"Sistem pakar membantu pemilik merpati mendiagnosis penyakit secara mandiri, cepat, dan akurat.\"",
            font_size=14, color=C_ACCENT, italic=True, font_name="Calibri", align=PP_ALIGN.CENTER)

add_page_number(slide, 15, dark=True)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(BASE_DIR, "output.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
